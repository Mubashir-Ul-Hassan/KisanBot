"""Generate KisanBot_Presentation.pptx — a styled project deck.

Run:  python scripts/make_ppt.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Brand palette (matches the app) ---------------------------------------
BG      = RGBColor(0x0B, 0x0F, 0x19)   # near-black navy
CARD    = RGBColor(0x15, 0x1B, 0x2B)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
RED     = RGBColor(0xEF, 0x44, 0x44)
LIGHT   = RGBColor(0xE5, 0xE7, 0xEB)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space=6.0):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = FONT
    return tb


def bullets(s, x, y, w, h, items, size=16, color=LIGHT, gap=10, marker="•  ",
            mcolor=EMERALD):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            head, body = it
        else:
            head, body = it, None
        rm = p.add_run(); rm.text = marker
        rm.font.size = Pt(size); rm.font.color.rgb = mcolor; rm.font.bold = True; rm.font.name = FONT
        rh = p.add_run(); rh.text = head
        rh.font.size = Pt(size); rh.font.color.rgb = color; rh.font.bold = bool(body); rh.font.name = FONT
        if body:
            rb = p.add_run(); rb.text = " — " + body
            rb.font.size = Pt(size); rb.font.color.rgb = MUTED; rb.font.name = FONT
    return tb


def header(s, kicker, title, accent=EMERALD):
    box(s, 0.6, 0.55, 0.14, 0.9, fill=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 0.9, 0.5, 11.8, 1.1, [
        [(kicker.upper(), 12, accent, True)],
        [(title, 30, WHITE, True)],
    ], space=2)


def pagefoot(s, n):
    text(s, 0.9, 7.0, 8, 0.4, [[("KisanBot — Multi-Agent Farm Advisory", 10, MUTED, False)]])
    text(s, 11.5, 7.0, 1.4, 0.4, [[(str(n), 10, MUTED, False)]], align=PP_ALIGN.RIGHT)


# ============================================================ SLIDE 1 — TITLE
s = slide()
box(s, 0, 0, SW, SH, fill=BG, shape=MSO_SHAPE.RECTANGLE)
box(s, 0.0, 0.0, 0.25, SH, fill=EMERALD, shape=MSO_SHAPE.RECTANGLE)
text(s, 1.0, 2.0, 11.5, 1.0, [[("🌾  KisanBot", 54, WHITE, True)]])
text(s, 1.05, 3.1, 11.5, 1.0, [[("A Collaborative Multi-Agent Farm Advisory System", 26, EMERALD, True)]])
text(s, 1.05, 4.0, 11.0, 1.6, [
    [("Urdu-speaking AI assistant for small-scale Pakistani farmers.", 18, LIGHT, False)],
    [("Powered by cooperating specialist agents — grounded in real data, honest about uncertainty.", 16, MUTED, False)],
])
box(s, 1.05, 5.7, 4.6, 0.6, fill=CARD, line=EMERALD)
text(s, 1.05, 5.72, 4.6, 0.56, [[("Assistant persona:  ساتھی  (Saathi)", 15, LIGHT, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ================================================= SLIDE 2 — WHO IT'S FOR
s = slide(); header(s, "The user", "Who this is for")
bullets(s, 0.9, 1.9, 11.6, 4.5, [
    ("Small-scale farmers across Pakistan", "the core user"),
    ("Limited formal education; some elderly", "many prefer speaking over reading/typing"),
    ("Basic Android phones", "on inconsistent, expensive mobile data"),
    ("Would rather just talk to it", "than navigate a cluttered app"),
], size=18, gap=16)
box(s, 0.9, 5.6, 11.6, 1.0, fill=CARD, line=BLUE)
text(s, 1.1, 5.62, 11.2, 0.96,
     [[("Design test: would this work for someone who finds a cluttered app confusing "
        "and would rather speak than type?", 15, LIGHT, False)]],
     anchor=MSO_ANCHOR.MIDDLE)
pagefoot(s, 2)

# ================================================= SLIDE 3 — THE PROBLEM (BEFORE)
s = slide(); header(s, "Starting point", "What the code was — a facade", accent=RED)
bullets(s, 0.9, 1.9, 11.8, 4.6, [
    ("\"Multi-agent\" was fake", "CrewAI / AutoGen / LangGraph / BeeAI imported but never used"),
    ("Each \"agent\" = one Gemini role-play prompt", "\"You are simulating a CrewAI crew…\""),
    ("Routing was a hardcoded if/elif chain", "not real tool/function calling"),
    ("Fabricated data", "NDVI/soil-moisture invented with random.seed; weather hardcoded by month"),
    ("No image diagnosis agent", "and pesticide advice had zero safety checks"),
    ("Two divergent brains", "client-side JS rules + Python server, duplicated logic"),
], size=17, gap=13, mcolor=RED)
pagefoot(s, 3)

# ================================================= SLIDE 4 — THE SOLUTION
s = slide(); header(s, "The solution", "A real cooperating agent team")
text(s, 0.9, 1.8, 11.8, 0.9,
     [[("One Orchestrator holds the conversation and routes to specialists. It never "
        "answers farming questions from its own knowledge — it delegates, then composes "
        "one clear Urdu reply.", 16, LIGHT, False)]])
cards = [
    ("🧭 Orchestrator", "Classify · ask 1 question · route · compose Urdu", EMERALD),
    ("🌤️ Weather", "Google Weather → farming implication", BLUE),
    ("🌱 Crop Recommendation", "Grounded in this location's soil + weather", EMERALD),
    ("🐛 Crop Protection", "Cultural-first · banned-list screen · disclaimer", AMBER),
    ("📷 Image Diagnosis", "Vision + explicit confidence → hands off", BLUE),
    ("💰 Market", "Approximate reference prices, clearly labelled", EMERALD),
]
x0, y0, w, h, gx, gy = 0.9, 2.9, 3.75, 1.55, 0.2, 0.25
for i, (t, d, c) in enumerate(cards):
    cx = x0 + (i % 3) * (w + gx)
    cy = y0 + (i // 3) * (h + gy)
    box(s, cx, cy, w, h, fill=CARD, line=c)
    text(s, cx + 0.15, cy + 0.12, w - 0.3, h - 0.2, [
        [(t, 15, WHITE, True)],
        [(d, 11.5, MUTED, False)],
    ], space=4)
pagefoot(s, 4)

# ================================================= SLIDE 5 — ARCHITECTURE FLOW
s = slide(); header(s, "How it works", "Orchestration flow")

def flow_box(x, y, w, h, label, sub, fill, line):
    box(s, x, y, w, h, fill=fill, line=line)
    text(s, x, y, w, h, [
        [(label, 14, WHITE, True)],
        [(sub, 10, MUTED, False)] if sub else [("", 10, MUTED, False)],
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=2)

flow_box(0.9, 2.2, 2.4, 1.1, "Farmer", "text · voice · photo", CARD, MUTED)
flow_box(4.0, 2.2, 2.7, 1.1, "Orchestrator", "classify + route", EMERALD, EMERALD)
flow_box(7.4, 1.1, 5.0, 0.85, "Ask ONE clarifying question", "if location/crop missing", CARD, BLUE)
flow_box(7.4, 2.35, 5.0, 0.85, "Specialist agent(s)", "grounded in real tools/data", CARD, EMERALD)
flow_box(7.4, 3.6, 5.0, 0.85, "Compose one Urdu reply", "+ disclaimer + confidence", CARD, AMBER)
flow_box(4.0, 4.9, 2.7, 1.0, "Farmer hears/reads", "simple Urdu answer", CARD, MUTED)

for a in [(3.3, 2.75, 4.0, 2.75), (6.7, 2.75, 7.4, 2.77), (6.7, 2.6, 7.4, 1.5),
          (9.9, 3.2, 9.9, 3.6), (7.4, 4.0, 5.35, 4.9)]:
    conn = s.shapes.add_connector(2, Inches(a[0]), Inches(a[1]), Inches(a[2]), Inches(a[3]))
    conn.line.color.rgb = EMERALD; conn.line.width = Pt(1.75)
text(s, 0.9, 6.2, 11.8, 0.6,
     [[("Lightweight LLM tool-calling — no heavy framework for 5–6 agents (per spec).", 13, MUTED, False)]])
pagefoot(s, 5)

# ================================================= SLIDE 6 — GROUNDING / TRUST
s = slide(); header(s, "The hard rule", "Non-hallucination & trust", accent=AMBER)
text(s, 0.9, 1.75, 11.8, 0.6, [[("A farmer's livelihood is at stake — so every answer is honest or absent.", 16, LIGHT, False)]])
bullets(s, 0.9, 2.5, 11.8, 4.2, [
    ("Never fabricate", "removed the random-seed NDVI and month-hardcoded weather"),
    ("Honest failures", "if data is unavailable, say so in Urdu and point to a human"),
    ("Ask, don't assume", "missing location/crop/symptom → one clarifying question"),
    ("State confidence", "low confidence + high stakes → 'confirm with your ag officer'"),
    ("Pesticide safety", "cultural controls first · banned list screened · dose never invented"),
    ("Traceability", "every answer logs which agent + which data source produced it"),
], size=17, gap=13, mcolor=AMBER)
pagefoot(s, 6)

# ================================================= SLIDE 7 — DATA & APIS
s = slide(); header(s, "Grounded in", "Data sources & APIs", accent=BLUE)
rows = [
    ("Gemini 2.5 Flash", "Agent intelligence + image diagnosis", "behind a swappable LLM interface"),
    ("Google Weather API", "Live conditions + forecast", "needs Google Maps key"),
    ("Google Geocoding", "Village/district → coordinates", "falls back to bundled regions"),
    ("ISRIC SoilGrids", "Real soil: pH, texture, carbon…", "pre-fetched & cached, never live per request"),
    ("Pesticide reference", "Banned-list screening", "curated; pending official DPP sign-off"),
    ("Browser Web Speech", "Urdu voice in/out (ur-PK)", "swappable for Google/Azure later"),
]
y = 1.95
for name, what, note in rows:
    box(s, 0.9, y, 11.6, 0.78, fill=CARD, line=None)
    text(s, 1.1, y, 3.3, 0.78, [[(name, 14, EMERALD, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.4, y, 4.0, 0.78, [[(what, 13, LIGHT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 8.5, y, 3.9, 0.78, [[(note, 12, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.86
pagefoot(s, 7)

# ================================================= SLIDE 8 — UX
s = slide(); header(s, "Built for the field", "Urdu-first, voice-first UX")
bullets(s, 0.9, 1.9, 6.0, 4.6, [
    ("Urdu-first by default", "no stray English shown to the farmer"),
    ("Voice in & out", "speak or type — your choice"),
    ("One big camera button", "photo → instant diagnosis"),
    ("One question at a time", "no deep menus"),
    ("Always-visible helpline", "reach a real human anytime"),
    ("Light & fault-tolerant", "friendly Urdu message when offline"),
], size=16, gap=13)
box(s, 7.2, 1.9, 5.3, 4.6, fill=CARD, line=EMERALD)
text(s, 7.45, 2.1, 4.9, 4.2, [
    [("Example", 13, EMERALD, True)],
    [("Farmer: my crop's leaves have spots", 13, LIGHT, False)],
    [("Saathi: don't worry — please send a clear "
      "photo of the affected leaf", 13, MUTED, False)],
    [("[photo → Image Diagnosis → Crop Protection]", 11, BLUE, True)],
    [("Saathi: this looks like leaf blight — I'm about "
      "80% confident. Please also confirm at your "
      "nearest agriculture center.", 13, LIGHT, False)],
], space=10)
pagefoot(s, 8)

# ================================================= SLIDE 9 — TECH / STRUCTURE
s = slide(); header(s, "Under the hood", "Tech stack & structure")
box(s, 0.9, 1.9, 5.7, 4.7, fill=CARD, line=None)
text(s, 1.15, 2.05, 5.3, 4.4, [
    [("Backend", 15, EMERALD, True)],
    [("Python · FastAPI (thin routing layer)", 13, LIGHT, False)],
    [("/agents  — orchestrator + 5 specialists", 12, MUTED, False)],
    [("/services — llm, geocoding, weather, soil,", 12, MUTED, False)],
    [("           pesticide, session, logging", 12, MUTED, False)],
    [("/agents/prompts — each prompt in its own file", 12, MUTED, False)],
    [("Endpoints: /api/chat · /api/diagnose · /health", 12, BLUE, False)],
], space=7)
box(s, 6.8, 1.9, 5.7, 4.7, fill=CARD, line=None)
text(s, 7.05, 2.05, 5.3, 4.4, [
    [("Frontend", 15, EMERALD, True)],
    [("Single-page app (HTML/CSS/JS)", 13, LIGHT, False)],
    [("Urdu-first · voice · photo upload · helpline", 12, MUTED, False)],
    [("Talks to the backend; offline rule fallback", 12, MUTED, False)],
    [("", 6, MUTED, False)],
    [("Config & safety", 15, EMERALD, True)],
    [("Keys via env only (.env) — never hardcoded", 12, MUTED, False)],
    [("Provider swappable (Gemini today → Claude ready)", 12, MUTED, False)],
], space=7)
pagefoot(s, 9)

# ================================================= SLIDE 10 — STATUS
s = slide(); header(s, "Status", "Done & verified live")
bullets(s, 0.9, 1.9, 11.8, 3.4, [
    ("Real multi-agent system replaces the facade", "orchestrator + 5 specialists"),
    ("Tested with a live Gemini key", "greeting, crop advice, pest diagnosis all working"),
    ("Crop advice is location-specific", "e.g. Multan kharif + canal water → cotton/maize/sugarcane"),
    ("Pest reply: cultural-first + safe pesticide + disclaimer", "banned substances stripped"),
    ("11/11 routing scenarios pass", "offline test needs no API key"),
], size=16, gap=12)
box(s, 0.9, 5.5, 11.6, 1.1, fill=CARD, line=EMERALD)
text(s, 1.1, 5.55, 11.2, 1.0,
     [[("Verified: python tests/scenarios.py → ALL ROUTING PASSED  ·  live Gemini "
        "end-to-end confirmed", 14, LIGHT, True)]], anchor=MSO_ANCHOR.MIDDLE)
pagefoot(s, 10)

# ================================================= SLIDE 11 — OPEN ITEMS
s = slide(); header(s, "Before real farmers", "Open items (developer sign-off)", accent=AMBER)
bullets(s, 0.9, 1.9, 11.8, 4.6, [
    ("Google Maps API key", "turns on live weather (everything else already works)"),
    ("Rotate the Gemini key", "regenerate a fresh one for security"),
    ("Verify pesticide list", "check data against the official DPP source"),
    ("Verify / localize helpline number", "per the user's province"),
    ("Populate soil cache", "run scripts/prefetch_soil.py once"),
    ("Native Urdu review", "have a speaker check the tone before shipping"),
], size=17, gap=13, mcolor=AMBER)
pagefoot(s, 11)

# ================================================= SLIDE 12 — CLOSING
s = slide()
box(s, 0.0, 0.0, 0.25, SH, fill=EMERALD, shape=MSO_SHAPE.RECTANGLE)
text(s, 1.0, 2.5, 11.5, 1.2, [[("Grounded. Honest. In Urdu.", 40, WHITE, True)]])
text(s, 1.05, 3.7, 11.3, 1.4, [
    [("A farm advisor that delegates to specialists, refuses to guess, and always "
      "offers a way to reach a real human.", 18, LIGHT, False)],
])
box(s, 1.05, 5.3, 6.2, 0.7, fill=CARD, line=EMERALD)
text(s, 1.05, 5.32, 6.2, 0.66, [[("Run:  python server.py  →  open index.html", 15, EMERALD, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "KisanBot_Presentation.pptx")
prs.save(out)
print("Saved:", out, "|", len(prs.slides._sldIdLst), "slides")
